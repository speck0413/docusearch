"""PPTX extraction (Phase 6 / GATE 6): slides → title heading, body, speaker notes, tables, images."""

from __future__ import annotations

import io
from pathlib import Path

from docusearch import ingest


def _build_pptx(path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    # Slide 1 — title + body + speaker notes
    s1 = prs.slides.add_slide(prs.slide_layouts[1])
    s1.shapes.title.text = "SPI Timing"
    s1.placeholders[1].text = "The SPI nonce ZQX7734 configures the peripheral bus."
    s1.notes_slide.notes_text_frame.text = "Speaker note: verify the strobe alignment before use."
    # Slide 2 — title + a table + a picture
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = "Registers"
    tbl = s2.shapes.add_table(2, 2, Inches(1), Inches(1.5), Inches(4), Inches(1)).table
    tbl.cell(0, 0).text = "Reg"
    tbl.cell(0, 1).text = "Val"
    tbl.cell(1, 0).text = "CTRL"
    tbl.cell(1, 1).text = "0x1"
    png = (
        b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x06\x00\x00"
        b"\x00\x1f\x15\xc4\x89\x00\x00\x00\nIDATx\x9cc\x00\x01\x00\x00\x05\x00\x01\r\n-\xb4"
        b"\x00\x00\x00\x00IEND\xaeB`\x82"
    )
    s2.shapes.add_picture(io.BytesIO(png), Inches(1), Inches(3), Inches(1), Inches(1))
    prs.save(str(path))


def test_extract_pptx_slides_notes_tables_images(tmp_path: Path) -> None:
    p = tmp_path / "deck.pptx"
    _build_pptx(p)
    doc = ingest.extract_pptx(p.read_bytes())

    bodies = "\n".join(s.text for s in doc.segments if s.kind == "body")
    assert "ZQX7734" in bodies  # slide body text
    assert "strobe alignment" in bodies  # speaker notes captured
    table = next(s for s in doc.segments if s.kind == "table")
    assert "Reg | Val" in table.text and "CTRL | 0x1" in table.text  # linearized
    assert doc.title  # a document title
    assert any("SPI Timing" in s.heading_path for s in doc.segments)  # slide title → heading path
    assert len(doc.images) == 1  # the picture is retained (R-ING-6)
    assert doc.images[0].data and doc.images[0].ext == "png"


def test_extract_document_dispatches_pptx(tmp_path: Path) -> None:
    p = tmp_path / "deck.pptx"
    _build_pptx(p)
    doc = ingest.extract_document(p, "pptx")
    assert any("ZQX7734" in s.text for s in doc.segments)

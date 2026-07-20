"""Multi-format report OUTPUT (Phase 6 / GATE 6).

Render a cited report to **any format docusearch can read, except STDF** — PDF, DOCX, PPTX, XLSX
(``md``/``html`` stay in :mod:`report`). Same citation discipline: a citation outside the evidence
set is refused (R-CIT-1) before a byte is written. Reuses the format libraries' write side
(python-docx / python-pptx / openpyxl / reportlab) — all pip, no LibreOffice.
"""

from __future__ import annotations

import io
import re
from collections.abc import Mapping, Sequence
from concurrent.futures import ThreadPoolExecutor
from contextlib import suppress
from pathlib import Path
from typing import Any

from . import citations

EXPORT_FORMATS = ("pdf", "docx", "pptx", "xlsx")

# plotly needs a tick to draw before the print snapshot is taken
_PDF_SETTLE_MS = 250


class ExportDependencyError(RuntimeError):
    """A format's writer is not installed. The message names the exact command to fix it."""

# Control chars that are illegal in OOXML/XML (keep tab, LF, CR) — a NUL in a title/body otherwise
# crashes python-docx/openpyxl with a raw traceback (red-team M1).
_CONTROL = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f]")


def xml_safe(text: str) -> str:
    """Strip control characters that OOXML writers reject."""
    return _CONTROL.sub("", text)


def xlsx_cell(value: str) -> str:
    """Sanitize a value for a spreadsheet cell: strip control chars AND neutralize **formula
    injection** — a leading ``= + - @`` is prefixed with ``'`` so the cell stays TEXT, never a live
    formula in whoever opens the file (CWE-1236, red-team H3)."""
    v = xml_safe(value)
    return "'" + v if v[:1] in ("=", "+", "-", "@") else v

# What a good report looks like in each format. The renderer can lay out whatever it is given,
# but it cannot invent structure that is not there: a section written as four dense paragraphs
# becomes four dense bullets on a slide. So the AUTHOR has to know the target before writing,
# and this is the single place that contract is stated — the MCP tool, help(), and the skill all
# read from here so they cannot drift.
FORMAT_GUIDANCE: dict[str, str] = {
    "md": (
        "PROSE DOCUMENT. Full paragraphs, fenced code, markdown lists, tables where they earn "
        "their place. No length limit — depth is the point."
    ),
    "html": (
        "PROSE DOCUMENT, the richest layout. Sections render as cards styled by their `kind`, so "
        "pick kinds deliberately (overview / procedure / code / hardware / config / test-program "
        "/ warning / reference). Figures render inline — pass the `img` shas of diagrams that "
        "genuinely explain something."
    ),
    "html-slide": (
        "A PRESENTATION in a browser — design it exactly like the pptx guidance below, including "
        "attaching each figure to the section it illustrates."
    ),
    "pptx": (
        "You are DESIGNING A PRESENTATION. Slides of uniform grey bullets are a bad deck no "
        "matter how accurate they are, and density is what makes them unreadable.\n"
        "  * ONE IDEA PER SECTION, and keep each section short — roughly 3-5 lines. A section "
        "that runs long is split across continuation slides, which reads worse than splitting it "
        "yourself into two ideas with their own headings.\n"
        "  * SHORT LINES. Aim for under ~12 words; lead with the point and cut the run-up. Long "
        "lines force the text smaller to fit, which is what makes a slide look cramped.\n"
        "  * ATTACH A FIGURE TO ITS SECTION (that section's `images`). A section with a figure "
        "and a few short points is laid out with the picture BESIDE the text — the best-looking "
        "slide this builder makes, and the one that actually explains something. Prefer it.\n"
        "  * GIVE IT AN ARC: what this is -> why it matters -> how to do it -> what goes wrong. "
        "Vary the sections so consecutive slides do not look identical.\n"
        "  * PUT THE DEPTH IN THE PROSE, not on the slide: the full section text is preserved in "
        "the speaker notes, so a short slide loses nothing. Write the slide for the eye and the "
        "notes for the presenter.\n"
        "  * Code: only the lines that matter, in its own section. A full listing does not fit a "
        "slide and never reads well on one."
    ),
    "xlsx": (
        "A GRID, one row per point. Each list item should be one self-contained fact; nested "
        "items ('  - ') become the detail column. It is meant to be sorted and filtered, so "
        "think in records, not paragraphs — a paragraph becomes one unreadable cell."
    ),
    "docx": (
        "A DOCUMENT for people to read and mark up. Prose with real headings, short paragraphs, "
        "markdown lists for steps and enumerations. Figures are welcome — pass the `img` shas of "
        "diagrams worth showing."
    ),
    "pdf": (
        "A DOCUMENT — authored exactly like docx; the PDF is printed from the HTML rendering, so "
        "figures and card styling carry through."
    ),
}

# What holds for every format. Deliberately short on rules: the only fixed parts of a report are
# the ones a reader needs to trust it.
_UNIVERSAL_GUIDANCE = (
    "FIXED (the builder adds these — do not write them yourself): the banner with its "
    "classification, request and provenance, and the References list built from your evidence.\n"
    "REQUIRED of you: every catalog claim carries its [D:doc#chunk], general knowledge carries "
    "[GK], and no claim goes beyond what the cited chunk says.\n"
    "FIGURES GO WITH THEIR SECTION. Put a diagram's sha in that section's own `images` list, not "
    "in the spec's top-level `images` — a figure belongs beside the text it explains, telling "
    "the story as it is told. Top-level `images` are for anything that genuinely belongs "
    "nowhere in particular, and they land at the end.\n"
    "EVERYTHING ELSE IS YOURS. How many sections, what they are called, how long they run, what "
    "order they take, which `kind` each one is, whether to lead with a summary or build to one, "
    "when a table beats prose, when a figure beats both. There is no required outline and no "
    "required length. Write the report this question deserves, and make it look like something "
    "a person would be glad to be handed.\n"
    "The visual theme follows the operator's configuration unless the requester asks for another "
    "— pass spec['theme'] when they do."
)

def guidance(fmt: str) -> str:
    """Authoring guidance for a target format (empty for an unknown one)."""
    specific = FORMAT_GUIDANCE.get(fmt.lower(), "")
    return f"{specific}\n\n{_UNIVERSAL_GUIDANCE}" if specific else ""


_AI_WARNING = "AI-generated — verify every claim against the cited sources before relying on it."


def _sections(
    sections: Sequence[Mapping[str, Any]] | None, body: str
) -> list[tuple[str, str, list[str]]]:
    """``(heading, body, image shas)``. A section's figures belong with it, not in a dump at the
    end where the reader has to work out what each one referred to."""
    if sections:
        out = []
        for sec in sections:
            raw = sec.get("images") or []
            shas = [str(x) for x in raw] if isinstance(raw, (list, tuple)) else []
            out.append((str(sec.get("heading", "")), str(sec.get("body", "")), shas))
        return out
    return [("", body, [])] if body else []


def _paragraphs(body: str) -> list[str]:
    return [xml_safe(ln.strip()) for ln in body.splitlines() if ln.strip()]


def export_report(
    *,
    title: str,
    sections: Sequence[Mapping[str, str]] | None = None,
    body: str = "",
    subtitle: str = "",
    evidence: set[tuple[int, int]],
    fmt: str,
    request: str = "",
    requested_by: str = "",
    model: str = "",
    classification: str = "Confidential",
    ref_targets: Mapping[tuple[int, int], tuple[str, str]] | None = None,
    html: str = "",  # fallback source for fmt="pdf"
    markdown: str = "",  # the rendered markdown - the preferred source for fmt="pdf"
    pptx_template: str = "",  # a .pptx/.potx whose theme + layouts the deck inherits
    figure_map: Mapping[str, tuple[str, str]] | None = None,  # sha -> (file path, caption)
) -> bytes:
    """Render the report to ``fmt`` bytes. Raises ``CitationError`` if any citation references a
    ``(doc_id, chunk_id)`` outside ``evidence`` (R-CIT-1)."""
    fmt = fmt.lower()
    if fmt not in EXPORT_FORMATS:
        raise ValueError(f"unknown export format {fmt!r}; expected one of {EXPORT_FORMATS} (or md/html via report.render_report)")
    secs = _sections(sections, body)
    surface = "\n".join([title, subtitle, *(b for _h, b, _i in secs)])
    violations = citations.verify(surface, evidence)
    if violations:
        raise citations.CitationError(
            f"report cites {len(violations)} chunk(s) outside its evidence set: "
            + ", ".join(f"D:{c.doc_id}#{c.chunk_id}" for c in violations[:10])
        )
    secs, refs = _numbered(title, subtitle, secs, ref_targets, evidence)
    meta = [x for x in (f"Request: {request}" if request else "",
                        f"For: {requested_by}" if requested_by else "",
                        f"Model: {model}" if model else "",
                        f"Classification: {classification}") if x]
    if fmt == "docx":
        return _to_docx(title, subtitle, secs, meta, refs, figure_map=figure_map)
    if fmt == "pptx":
        return _to_pptx(title, subtitle, secs, meta, refs, template=pptx_template,
                        figure_map=figure_map)
    if fmt == "xlsx":
        return _to_xlsx(title, secs, meta, refs)
    # PDF is a DOCUMENT. It is built from the MARKDOWN rendering — a plain document flow —
    # rather than the web layout, whose cards and full-bleed banner print like a saved web page.
    if markdown:
        return html_to_pdf(markdown_to_print_html(markdown, title))
    if html:
        return html_to_pdf(html)
    raise ValueError("pdf export needs the rendered markdown - pass markdown=render_report(...)")


def _numbered(
    title: str,
    subtitle: str,
    secs: list[tuple[str, str, list[str]]],
    ref_targets: Mapping[tuple[int, int], tuple[str, str]] | None,
    evidence: set[tuple[int, int]],
    base_url: str = "",
) -> tuple[list[tuple[str, str, list[str]]], list[str]]:
    """Bodies with inline citations turned into ``[1]`` markers, plus the numbered reference list.

    Reuses the HTML/Markdown renderer's numbering so all six formats number identically
    (R-REUSE-2). A reference names the real document — ``store - title - heading`` — and never a
    bare ``D:<doc>#<chunk>``: that id is an internal key, and a report that shows one, or whose
    references are ids instead of documents, is not acceptable output.
    """
    from . import report

    numbering, ordered = report._collect_numbering([title, subtitle, *(b for _h, b, _i in secs)])
    bodies = [(head, report._cite_md(body, numbering), imgs) for head, body, imgs in secs]
    refs = [label for _href, label in report._references(ordered, base_url, ref_targets)]
    # evidence that was supplied but never cited still belongs in the list, after the cited ones
    cited = {(c.doc_id, c.chunk_id) for c in ordered}
    for key in sorted(evidence - cited):
        if ref_targets and key in ref_targets:
            refs.append(ref_targets[key][1])
    return bodies, [xml_safe(r) for r in refs]



# A slide holds a handful of points before it becomes a wall of text. Past this the section
# continues on another slide rather than overflowing off the bottom.
_SLIDE_BULLETS = 6
_SLIDE_CHARS = 160  # a "point" longer than this is prose, not a bullet
_SOURCE_LINES = 10  # a source list is one line each, so it packs tighter than bullets
_DIVIDER_MIN_SECTIONS = 8  # below this a deck is short enough that dividers just pad it


def _points(body: str) -> list[tuple[int, str]]:
    """A body as ``(indent level, text)`` bullet points.

    Markdown list markers become real bullet levels; a fenced code block keeps its lines
    verbatim. Prose that is not a list still yields one point per paragraph, so a document-shaped
    section degrades to readable bullets instead of one unbroken blob."""
    points: list[tuple[int, str]] = []
    in_code = False
    for raw in body.splitlines():
        line = raw.rstrip()
        if line.strip().startswith("```"):
            in_code = not in_code
            continue
        if not line.strip():
            continue
        if in_code:
            points.append((1, xml_safe(line.strip())))
            continue
        indent = (len(line) - len(line.lstrip())) // 2
        text = line.strip()
        if text[:2] in ("- ", "* ") or (text[:1].isdigit() and text[1:3] in (". ", ") ")):
            text = text.split(" ", 1)[1] if " " in text else text
            points.append((min(indent, 4), xml_safe(text)))
        else:
            points.append((0, xml_safe(text)))
    return points


def _slide_chunks(
    points: list[tuple[int, str]], per_slide: int = _SLIDE_BULLETS
) -> list[list[tuple[int, str]]]:
    """Split points across slides, keeping long prose points on lighter slides."""
    chunks: list[list[tuple[int, str]]] = []
    current: list[tuple[int, str]] = []
    budget = 0
    for level, text in points:
        cost = 2 if len(text) > _SLIDE_CHARS else 1
        if current and budget + cost > per_slide:
            chunks.append(current)
            current, budget = [], 0
        current.append((level, text))
        budget += cost
    if current:
        chunks.append(current)
    return chunks or [[]]


def _fill(holder: Any, points: list[tuple[int, str]]) -> None:
    """Write real bullet PARAGRAPHS with breathing room.

    Setting ``.text`` with newlines produces a single paragraph that neither indents nor bullets
    and simply overflows the placeholder. Line spacing and space-between-points are what separate
    a deck that looks composed from one that looks dumped: default OOXML paragraphs sit flush
    against each other, which reads as a wall even when the wording is short."""
    from pptx.util import Pt

    frame = holder.text_frame
    frame.clear()
    frame.word_wrap = True
    for i, (level, text) in enumerate(points or [(0, " ")]):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.text = text
        para.level = level
        para.line_spacing = 1.15
        para.space_after = Pt(10 if level == 0 else 4)


def _styled(paragraph: object, *names: str) -> None:
    """Apply the first BUILT-IN style that exists in this document.

    Reports are styled by NAME only — no hardcoded fonts, sizes or colours — so a corporate
    .dotx or theme restyles the whole report on drop-in. A template need not define every
    built-in and python-docx raises on an unknown style, so fall through and finally leave the
    default rather than failing an export over cosmetics."""
    for name in names:
        try:
            paragraph.style = name  # type: ignore[attr-defined]
            return
        except (KeyError, ValueError):
            continue


def _to_docx(
    title: str, subtitle: str, secs: list[tuple[str, str, list[str]]], meta: list[str], refs: list[str],
    *, figure_map: Mapping[str, tuple[str, str]] | None = None,
) -> bytes:
    from docx import Document
    from docx.shared import Inches

    figure_map = dict(figure_map or {})
    doc = Document()
    doc.add_heading(xml_safe(title), 0)
    if subtitle:
        _styled(doc.add_paragraph(xml_safe(subtitle)), "Subtitle", "Body Text")
    _styled(doc.add_paragraph(_AI_WARNING), "Intense Quote", "Quote", "Body Text")
    for line in meta:
        _styled(doc.add_paragraph(xml_safe(line)), "Caption", "Body Text")
    placed: set[str] = set()
    for heading, bdy, sec_imgs in secs:
        if heading:
            doc.add_heading(xml_safe(heading), level=1)
        for para in _paragraphs(bdy):
            _styled(doc.add_paragraph(para), "Body Text")
        for sha in sec_imgs:  # this section's figures, beside the text they explain
            if sha in figure_map:
                path, caption = figure_map[sha]
                with suppress(Exception):  # a bad image costs its figure, never the report
                    doc.add_picture(path, width=Inches(6.0))
                if caption:
                    _styled(doc.add_paragraph(caption), "Caption", "Body Text")
                placed.add(sha)
    loose = [(p, c) for sha, (p, c) in figure_map.items() if sha not in placed]
    if loose:  # anything the author did not attach to a section
        doc.add_heading("Figures", level=1)
        for path, caption in loose:
            with suppress(Exception):
                doc.add_picture(path, width=Inches(6.0))
            if caption:
                _styled(doc.add_paragraph(caption), "Caption", "Body Text")
    if refs:
        doc.add_heading("References", level=1)
        for i, r in enumerate(refs, 1):
            _styled(doc.add_paragraph(f"{i}. {r}"), "Body Text")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _layout(prs: Any, *names: str, fallback: int) -> Any:
    """A slide layout by NAME, falling back to an index.

    Every template orders its layouts differently, so taking ``slide_layouts[1]`` and hoping is
    how a themed deck comes out wrong. Match the standard names first."""
    # Preference order is the CALLER's, not the template's: scanning the deck's layouts and
    # taking the first match returned "Title Only" for a picture slide simply because it is
    # earlier in the file than "Picture with Caption".
    by_name = {layout.name.lower(): layout for layout in prs.slide_layouts}
    for name in names:
        found = by_name.get(name.lower())
        if found is not None:
            return found
    return prs.slide_layouts[fallback if fallback < len(prs.slide_layouts) else 0]


def _body_placeholder(slide: Any) -> Any:
    """The slide's content placeholder, found by TYPE not index.

    ``placeholders[1]`` is the body only by convention; in a real template that index may be a
    picture, a date, or absent entirely."""
    from pptx.enum.shapes import PP_PLACEHOLDER

    # Exclude the title by placeholder TYPE, not by object identity: python-pptx hands back a new
    # wrapper each time you touch `shapes.title`, so `p != slide.shapes.title` can be true for the
    # title itself — and the title would then be filled with body text.
    titles = (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
    holders = [p for p in slide.placeholders if p.placeholder_format.type not in titles]
    for kind in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.SUBTITLE):
        for holder in holders:
            if holder.placeholder_format.type == kind:
                return holder
    return holders[0] if holders else None


def _figure_slide(prs: Any, layout: Any, path: str, caption: str) -> None:
    """A slide for one figure: picture in the picture placeholder, caption beneath it."""
    from pptx.enum.shapes import PP_PLACEHOLDER

    slide = prs.slides.add_slide(layout)
    if slide.shapes.title is not None:
        slide.shapes.title.text = xml_safe(caption) or "Figure"
    holders = _content_holders(slide)
    picture = next(
        (h for h in holders if h.placeholder_format.type == PP_PLACEHOLDER.PICTURE), None
    )
    if picture is not None:
        _fill_box(prs, slide, picture, path)
        for other in holders:
            if other is not picture and not other.text_frame.text.strip():
                other._element.getparent().remove(other._element)
    else:
        for other in holders:
            other._element.getparent().remove(other._element)
        _place_picture(prs, slide, path)


def _place_picture(prs: Any, slide: Any, path: str) -> None:
    """Centre a picture in the slide's usable area, scaled to fit.

    Geometry comes from the template's own slide size, so a 16:9 corporate deck and the default
    4:3 both place it correctly — this is the one place positioning is unavoidable, and it is
    derived, never hardcoded."""
    from PIL import Image  # pillow ships with the vision extras; guarded by the caller

    top_margin = int(prs.slide_height * 0.22)  # below the title
    max_w = int(prs.slide_width * 0.86)
    max_h = int(prs.slide_height - top_margin - prs.slide_height * 0.08)
    try:
        with Image.open(path) as img:
            ratio = img.height / img.width if img.width else 0.75
    except Exception:  # noqa: BLE001 - unreadable image: fall back to a sane box
        ratio = 0.75
    width = max_w
    height = int(width * ratio)
    if height > max_h:
        height = max_h
        width = int(height / ratio) if ratio else max_w
    left = int((prs.slide_width - width) / 2)
    # A corrupt or unreadable image must cost its slide, never the whole report.
    with suppress(Exception):
        slide.shapes.add_picture(path, left, top_margin, width=width, height=height)


def _fit(frame: Any, points: list[tuple[int, str]]) -> None:
    """Shrink dense text so it cannot overflow its placeholder.

    Only ever shrinks: a template's own sizing is respected until the content would spill, which
    is the failure that makes a generated deck look broken. python-pptx cannot measure text, so
    this is a density heuristic, not exact layout."""
    from pptx.util import Pt

    chars = sum(len(t) for _l, t in points)
    lines = len(points)
    if chars <= 260 and lines <= 5:
        return  # comfortable — leave the template alone
    size = 18 if (chars <= 420 and lines <= 7) else (16 if chars <= 620 else 14)
    for para in frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(size)


def _fill_box(prs: Any, slide: Any, holder: Any, path: str) -> None:
    """Put a picture where a placeholder is, scaled to fit its box, then drop the placeholder."""
    from PIL import Image

    left, top = holder.left, holder.top
    box_w, box_h = holder.width, holder.height
    holder._element.getparent().remove(holder._element)
    try:
        with Image.open(path) as img:
            ratio = img.height / img.width if img.width else 0.75
    except Exception:  # noqa: BLE001
        ratio = 0.75
    width = box_w
    height = int(width * ratio)
    if height > box_h:
        height = box_h
        width = int(height / ratio) if ratio else box_w
    with suppress(Exception):
        slide.shapes.add_picture(
            path, left + int((box_w - width) / 2), top + int((box_h - height) / 2),
            width=width, height=height,
        )


def _content_holders(slide: Any) -> list[Any]:
    """The non-title placeholders of a slide, in layout order."""
    from pptx.enum.shapes import PP_PLACEHOLDER

    skip = (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.DATE,
            PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.SLIDE_NUMBER)
    return [p for p in slide.placeholders if p.placeholder_format.type not in skip]


def _to_pptx(
    title: str, subtitle: str, secs: list[tuple[str, str, list[str]]], meta: list[str],
    refs: list[str],
    *, template: str = "", figure_map: Mapping[str, tuple[str, str]] | None = None,
) -> bytes:
    """Build the deck from the template's LAYOUTS and PLACEHOLDERS.

    Uses more than one layout on purpose: a deck where every slide is a title over a bullet list
    is exhausting to sit through. A section with a figure puts the picture BESIDE its points
    (Two Content) so the two read together; a further figure gets a Picture-with-Caption slide;
    dense text is shrunk rather than allowed to overflow. Nothing sets a font face or colour, so
    swapping the template still restyles everything."""
    from pptx import Presentation

    figure_map = dict(figure_map or {})
    prs = Presentation(template) if template else Presentation()
    cover = prs.slides.add_slide(_layout(prs, "Title Slide", fallback=0))
    if cover.shapes.title is not None:
        cover.shapes.title.text = xml_safe(title)
    holders = _content_holders(cover)
    if holders:
        # subtitle then the provenance line, so the cover states what this is and who it is for
        sub = [(0, xml_safe(subtitle))] if subtitle else []
        sub += [(0, xml_safe(line)) for line in meta if line.startswith(("For:", "Request:"))]
        _fill(holders[0], sub or [(0, _AI_WARNING)])
        for extra in holders[1:]:
            extra._element.getparent().remove(extra._element)

    body_layout = _layout(prs, "Title and Content", "Title and Body", fallback=1)
    two_layout = _layout(prs, "Two Content", "Comparison", fallback=1)
    pic_layout = _layout(prs, "Picture with Caption", "Title Only", fallback=1)

    section_layout = _layout(prs, "Section Header", "Title Only", fallback=1)
    placed: set[str] = set()
    # A divider every few sections gives a long deck rhythm; on a short one it is just noise.
    divide_every = 4 if len(secs) >= _DIVIDER_MIN_SECTIONS else 0
    for idx, (heading, bdy, sec_imgs) in enumerate(secs):
        if divide_every and idx and idx % divide_every == 0:
            divider = prs.slides.add_slide(section_layout)
            if divider.shapes.title is not None:
                divider.shapes.title.text = xml_safe(heading)
            for extra in _content_holders(divider):
                extra._element.getparent().remove(extra._element)
        points = _points(bdy)
        figs = [(sha, figure_map[sha]) for sha in sec_imgs if sha in figure_map]
        chunks = _slide_chunks(points)
        head = xml_safe(heading) or xml_safe(title)
        for n, chunk in enumerate(chunks):
            pair = figs[0] if (n == 0 and figs and len(chunk) <= _SLIDE_BULLETS) else None
            layout = two_layout if pair is not None else body_layout
            slide = prs.slides.add_slide(layout)
            if slide.shapes.title is not None:
                slide.shapes.title.text = head if n == 0 else f"{head} (cont.)"
            holders = _content_holders(slide)
            if holders:
                _fill(holders[0], chunk)
                _fit(holders[0].text_frame, chunk)
            if pair is not None:
                sha, (path, _caption) = pair
                if len(holders) > 1:
                    _fill_box(prs, slide, holders[1], path)  # picture beside the points
                else:
                    _place_picture(prs, slide, path)
                placed.add(sha)
            for extra in holders[2:]:  # unused placeholders would show prompt text
                extra._element.getparent().remove(extra._element)
            if n == 0 and bdy.strip():
                with suppress(AttributeError, KeyError):  # a template without a notes master
                    slide.notes_slide.notes_text_frame.text = xml_safe(bdy.strip())
        for sha, (path, caption) in figs:  # any further figure gets its own slide
            if sha not in placed:
                _figure_slide(prs, pic_layout, path, caption)
                placed.add(sha)

    for sha, (path, caption) in figure_map.items():  # not attached to any section
        if sha not in placed:
            _figure_slide(prs, pic_layout, path, caption)

    if refs:
        # A deck that ends in eight slides of numbered chunk references is unreadable, but the
        # numbering has to survive because the inline markers point at it. So the SLIDE lists the
        # distinct source documents — what an audience can actually use — and the complete
        # numbered list goes in the speaker notes, where it stays auditable and matches the
        # inline numbers exactly.
        seen: dict[str, None] = {}
        for ref in refs:
            seen.setdefault(" — ".join(ref.split(" — ")[:2]).strip(), None)
        documents = [(0, name) for name in seen]
        full = "\n".join(f"{i}. {r}" for i, r in enumerate(refs, 1))
        for n, chunk in enumerate(_slide_chunks(documents, _SOURCE_LINES)):
            slide = prs.slides.add_slide(body_layout)
            if slide.shapes.title is not None:
                label = f"Sources ({len(seen)} documents, {len(refs)} citations)"
                slide.shapes.title.text = label if n == 0 else "Sources (cont.)"
            holders = _content_holders(slide)
            if holders:
                _fill(holders[0], chunk)
                _fit(holders[0].text_frame, chunk)
            if n == 0:
                with suppress(AttributeError, KeyError):
                    slide.notes_slide.notes_text_frame.text = full
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _to_xlsx(title: str, secs: list[tuple[str, str, list[str]]], meta: list[str], refs: list[str]) -> bytes:
    """A spreadsheet is a GRID, so emit rows and columns — not a document poured down column A.

    Section / Point / Detail lets you filter and sort by section, which is the only reason to
    want a report as a workbook. Provenance and references get their own sheets rather than
    interrupting the data."""
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font

    wb = Workbook()
    ws = wb.active
    ws.title = "Report"
    ws.append(["Section", "Point", "Detail"])
    for cell in ws[1]:
        cell.font = Font(bold=True)
    for heading, bdy, _imgs in secs:
        points = _points(bdy)
        for level, text in points:
            ws.append([
                xlsx_cell(heading),
                xlsx_cell(text) if level == 0 else "",
                xlsx_cell(text) if level > 0 else "",
            ])
    ws.freeze_panes = "A2"  # the header stays put while you scroll
    ws.auto_filter.ref = ws.dimensions  # filter/sort by section out of the box
    for col, width in (("A", 28), ("B", 90), ("C", 90)):
        ws.column_dimensions[col].width = width
    for row in ws.iter_rows(min_row=2):
        for cell in row:
            cell.alignment = Alignment(vertical="top", wrap_text=True)

    info = wb.create_sheet("About")
    info.append(["Title", xlsx_cell(title)])
    info.append(["Notice", _AI_WARNING])
    for line in meta:
        key, _, value = line.partition(": ")
        info.append([xlsx_cell(key), xlsx_cell(value)])
    info.column_dimensions["A"].width = 18
    info.column_dimensions["B"].width = 90

    if refs:
        sheet = wb.create_sheet("References")
        sheet.append(["#", "Source"])
        for cell in sheet[1]:
            cell.font = Font(bold=True)
        for i, r in enumerate(refs, 1):
            sheet.append([i, xlsx_cell(r)])
        sheet.freeze_panes = "A2"
        sheet.column_dimensions["A"].width = 5
        sheet.column_dimensions["B"].width = 110

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _soffice() -> str | None:
    """The LibreOffice binary, if this machine has one."""
    import shutil

    for name in ("soffice", "libreoffice"):
        found = shutil.which(name)
        if found:
            return found
    mac = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    return mac if Path(mac).is_file() else None


_PRINT_CSS = """
@page{size:Letter;margin:0.85in 0.8in;}
body{font:11.5pt/1.55 Georgia,"Times New Roman",serif;color:#111;max-width:none;margin:0;}
h1{font-size:20pt;margin:0 0 .2em;line-height:1.2;}
h2{font-size:14.5pt;margin:1.4em 0 .35em;border-bottom:1px solid #ccc;padding-bottom:.15em;
page-break-after:avoid;}
h3{font-size:12.5pt;margin:1.1em 0 .3em;page-break-after:avoid;}
p,li{orphans:3;widows:3;}
blockquote{margin:.6em 0;padding:.5em .9em;background:#f4f4f6;border-left:3px solid #bbb;
font-size:10.5pt;color:#333;}
pre{background:#f4f4f6;border:1px solid #ddd;border-radius:4px;padding:.7em .9em;
font:9.5pt/1.4 "SF Mono",Menlo,Consolas,monospace;white-space:pre-wrap;page-break-inside:avoid;}
code{font:9.5pt "SF Mono",Menlo,Consolas,monospace;background:#f4f4f6;padding:.1em .3em;
border-radius:3px;}
pre code{background:none;padding:0;}
img{max-width:100%;max-height:4.2in;height:auto;display:block;margin:.7em auto;
page-break-inside:avoid;}
table{border-collapse:collapse;width:100%;font-size:10pt;margin:.7em 0;}
th,td{border:1px solid #ccc;padding:.35em .5em;text-align:left;}
th{background:#f0f0f3;}
hr{border:none;border-top:1px solid #ccc;margin:1.4em 0;}
a{color:#123f7a;}
"""


def markdown_to_print_html(md: str, title: str = "") -> str:
    """Markdown -> a document-shaped HTML page for printing.

    Deliberately not the report's web layout: a PDF should read as a document, so this is a
    single serif column with print-aware page breaks and no cards or banner styling."""
    from html import escape

    try:
        from markdown_it import MarkdownIt

        body = MarkdownIt("commonmark", {"html": False, "linkify": False}).enable("table").render(md)
    except ImportError:  # keep working without the md extra
        body = "<pre>" + escape(md) + "</pre>"
    return (
        '<!doctype html><html lang="en"><head><meta charset="utf-8">'
        f"<title>{escape(title or 'Report')}</title><style>{_PRINT_CSS}</style>"
        f"</head><body>{body}</body></html>"
    )


def docx_to_pdf(data: bytes) -> bytes | None:
    """Convert a .docx to PDF with LibreOffice, or None when no converter is installed.

    Returning None rather than raising is deliberate: PDF must keep working on a machine without
    LibreOffice, so the caller falls back to printing the HTML. The docx route is preferred
    because a PDF should look like a document, not like a saved web page."""
    binary = _soffice()
    if binary is None:
        return None
    import subprocess
    import tempfile

    with tempfile.TemporaryDirectory() as tmp:
        src = Path(tmp) / "report.docx"
        src.write_bytes(data)
        try:
            subprocess.run(  # noqa: S603 - argv built here, paths are ours
                [binary, "--headless", "--norestore", "--convert-to", "pdf",
                 "--outdir", tmp, str(src)],
                capture_output=True, timeout=180, check=False,
            )
        except (OSError, subprocess.SubprocessError):
            return None
        out = Path(tmp) / "report.pdf"
        return out.read_bytes() if out.is_file() else None


def html_to_pdf(html: str) -> bytes:
    """Print a rendered HTML report to PDF in headless chromium.

    PDF is the one output format with no structured writer of its own, deliberately: rendering
    the real HTML means there is ONE layout to maintain instead of two that drift, and because
    chromium executes the page's JS, interactive plotly charts appear in the PDF. A direct
    writer (reportlab) or a JS-less converter (WeasyPrint) would silently drop them.
    """
    try:
        from playwright.sync_api import sync_playwright
    except ImportError as err:  # actionable, one line (operability contract)
        raise ExportDependencyError(
            "PDF export needs the [export] extra: pip install 'docusearch[export]' "
            "&& python -m playwright install chromium"
        ) from err
    def render() -> bytes:
        with sync_playwright() as pw:
            browser = pw.chromium.launch()
            try:
                page = browser.new_page()
                # The report is self-contained (inlined CSS/JS, data: images) so there is no
                # network to wait on, but plotly needs a tick to draw before the print snapshot.
                page.set_content(html, wait_until="domcontentloaded")
                page.wait_for_timeout(_PDF_SETTLE_MS)
                return bytes(
                    page.pdf(
                        format="Letter",
                        print_background=True,  # keep the report's banner/card styling
                        margin={"top": "0.6in", "bottom": "0.6in",
                                "left": "0.5in", "right": "0.5in"},
                    )
                )
            finally:
                browser.close()

    try:
        # Playwright's SYNC api refuses to run inside a running asyncio loop, and the server is
        # async — calling it on the request thread fails with a bare "Error" that says nothing.
        # A worker thread has no loop of its own, so this works from both the async server and a
        # plain synchronous caller like the CLI.
        with ThreadPoolExecutor(max_workers=1) as pool:
            return pool.submit(render).result()
    except ExportDependencyError:
        raise
    except Exception as err:  # a missing browser download is the common case
        raise ExportDependencyError(
            f"PDF export could not drive chromium ({type(err).__name__}). Run: "
            "python -m playwright install chromium"
        ) from err

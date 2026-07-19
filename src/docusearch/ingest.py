"""Ingestion pipeline: filesystem -> extract -> chunk -> link/image -> index (§7).

The largest module by design (R-ARCH-3). It turns a source folder of documents into
rows in the store: discover files (globs), skip unchanged ones by content hash
(R-ING-3), strip boilerplate and extract structured text (R-ING-2, §7.3), chunk while
preserving code blocks (R-ING-4), capture links and images (R-ING-5/6), index into FTS5,
and emit a loud audit report (§7.8).

Public surface (grows through Phase 1):
    iter_files(location, include, exclude) -> Iterator[Path]   # source discovery
    content_hash(path) -> str                                  # SHA-256, incremental skip
    extract_html(html, *, content_selector, strip_selectors) -> ExtractedDoc  # §7.3
    ExtractedDoc / Segment / LinkRef / ImageRef               # extraction result types
    chunk_document(doc, *, chunk_tokens, overlap) -> list[Chunk]   # §7.6
    Chunk                                                     # a ready-to-index chunk
    run_ingest(config, store, *, force) -> IngestResult       # the whole pipeline (§7)
    IngestResult                                              # per-run audit counts
    render_ingest_audit(result, *, run_id) -> str             # §7.8 report
    render_store_audit(store) -> str                          # `docusearch audit`
"""

from __future__ import annotations

import contextlib
import hashlib
import io
import os
import re
import time
from collections.abc import Callable, Iterator, Sequence
from dataclasses import dataclass, field
from pathlib import Path

from selectolax.parser import HTMLParser, Node

from . import embed, runlog
from .config import Config, SourceConfig
from .embed import EmbedProvider
from .enrich import GotchaPattern, active_gotcha_patterns, gotcha_tag_text, match_gotcha
from .store import Store

# A progress sink: (phase, done, total). ``phase`` is "ingest" (files) or "embed"
# (chunks). Optional everywhere — library callers pass nothing; the CLI renders a bar.
ProgressFn = Callable[[str, int, int], None]

_HASH_BLOCK = 1 << 20

_HEADINGS = {"h1", "h2", "h3", "h4", "h5", "h6"}
# Never indexed and never mined for relations: scripting/markup noise plus semantic navigation
# chrome (nav/aside/footer). A cross-reference counts only when it lives in the BODY text we index —
# a link in the navigation pane does not (Stephen, 2026-07-18). Div/class-based nav (e.g. a Sphinx
# sidebar) is still handled by `content_selector` / `strip_selectors`, since it carries no semantic tag.
_SKIP = {"script", "style", "template", "svg", "math", "noscript", "nav", "aside", "footer"}
# Block-level tags: an element with NO block-level element child is a "leaf block" whose
# text is emitted whole (with a separator so inline children never glue); an element WITH
# a block child is a container we recurse into (buffering its own inline/text runs).
# Everything not listed here is treated as inline. pre/table/figure/img are special-cased.
_BLOCK = _HEADINGS | {
    "html",
    "body",
    "pre",
    "table",
    "figure",
    "img",
    "div",
    "section",
    "article",
    "main",
    "header",
    "footer",
    "aside",
    "nav",
    "p",
    "ul",
    "ol",
    "dl",
    "li",
    "dd",
    "dt",
    "blockquote",
    "form",
    "fieldset",
    "details",
    "summary",
    "address",
    "hr",
    "figcaption",
    "tbody",
    "thead",
    "tfoot",
    "tr",
    "td",
    "th",
    "caption",
}


def _glob_to_regex(pattern: str) -> re.Pattern[str]:
    """Translate a path glob to a segment-aware regex.

    ``**`` spans any number of path segments, ``*`` matches within one segment, ``?``
    matches one non-separator character. Matching is done against a POSIX relative path.
    """
    out: list[str] = []
    i, n = 0, len(pattern)
    while i < n:
        if pattern[i : i + 3] == "**/":
            out.append("(?:.*/)?")
            i += 3
        elif pattern[i : i + 2] == "**":
            out.append(".*")
            i += 2
        elif pattern[i] == "*":
            out.append("[^/]*")
            i += 1
        elif pattern[i] == "?":
            out.append("[^/]")
            i += 1
        else:
            out.append(re.escape(pattern[i]))
            i += 1
    return re.compile("^" + "".join(out) + "$")


def _classify_files(
    location: Path | str,
    include: Sequence[str],
    exclude: Sequence[str],
) -> tuple[list[Path], int, int]:
    """Sort files under ``location`` into (included, excluded_by_glob, other) for the audit.

    A file matching any exclude glob is excluded; otherwise it is included if it matches
    any include glob (an empty include list matches everything, R-ING-1); anything left
    is "other" (present but not selected). Included files come back deterministically
    sorted for reproducible ingest order.
    """
    root = Path(location)
    inc = [_glob_to_regex(p) for p in include] or [re.compile(".*")]
    exc = [_glob_to_regex(p) for p in exclude]
    included: list[Path] = []
    excluded_glob = 0
    other = 0
    for path in sorted(root.rglob("*")):
        if not path.is_file():
            continue
        rel = path.relative_to(root).as_posix()
        if any(r.match(rel) for r in exc):
            excluded_glob += 1
        elif any(r.match(rel) for r in inc):
            included.append(path)
        else:
            other += 1
    return included, excluded_glob, other


def iter_files(
    location: Path | str,
    include: Sequence[str],
    exclude: Sequence[str],
) -> Iterator[Path]:
    """Yield files under ``location`` matching any include glob and no exclude glob."""
    return iter(_classify_files(location, include, exclude)[0])


def content_hash(path: Path | str) -> str:
    """SHA-256 of a file's bytes — the incremental-ingest key (R-ING-3)."""
    h = hashlib.sha256()
    with Path(path).open("rb") as fh:
        for block in iter(lambda: fh.read(_HASH_BLOCK), b""):
            h.update(block)
    return h.hexdigest()


# ----------------------------------------------------------------- HTML extraction


@dataclass
class Segment:
    """One ordered piece of a document's content, tagged with its heading path."""

    kind: str  # body | code | table
    text: str
    heading_path: str


@dataclass
class LinkRef:
    """A cross-reference as written in the source (resolved later, R-ING-5)."""

    target: str  # raw href
    anchor: str
    link_type: str = "html_href"


@dataclass
class ImageRef:
    """An image reference kept for retention + text findability (R-ING-6).

    HTML images are referenced by ``src`` (resolved to a file on disk at retention time). PDF/DOCX
    images are **embedded inside the document**, so their bytes are carried inline via ``data``
    (+ ``ext``) and staged directly — no external file to resolve."""

    src: str
    alt: str
    caption: str
    heading_path: str
    data: bytes | None = None  # inline image bytes (PDF/DOCX); None => resolve ``src`` from disk
    ext: str = ""  # inline image extension (e.g. "png", "jpeg") when ``data`` is set


@dataclass
class ExtractedDoc:
    title: str
    segments: list[Segment] = field(default_factory=list)
    links: list[LinkRef] = field(default_factory=list)
    images: list[ImageRef] = field(default_factory=list)
    content_selector_matched: bool = True

    @property
    def text_length(self) -> int:
        """Total visible text length — compared against ``min_content_chars`` (R-ING-2)."""
        return sum(len(s.text) for s in self.segments)


def _clean(text: str) -> str:
    """Collapse runs of whitespace (for prose/tables/headings, never for code)."""
    return " ".join(text.split())


def _heading_path(stack: list[tuple[int, str]]) -> str:
    return " > ".join(text for _, text in stack)


def _push_heading(stack: list[tuple[int, str]], level: int, text: str) -> None:
    while stack and stack[-1][0] >= level:
        stack.pop()
    stack.append((level, text))


def _linearize_table(node: Node) -> str:
    rows: list[str] = []
    caption = node.css_first("caption")
    if caption is not None:
        cap = _clean(caption.text(separator=" "))
        if cap:
            rows.append(cap)
    for tr in node.css("tr"):
        cells = [_clean(cell.text(separator=" ")) for cell in tr.css("td, th")]
        if any(cells):
            rows.append(" | ".join(cells))
    return "\n".join(rows)


def _add_link(node: Node, doc: ExtractedDoc) -> None:
    href = node.attributes.get("href")
    if href:
        doc.links.append(LinkRef(target=href, anchor=_clean(node.text())))


def _add_image(node: Node, heading_path: str, doc: ExtractedDoc, caption: str = "") -> None:
    src = node.attributes.get("src")
    if not src:
        return
    # An inline data: URI carries the image bytes in the src itself — decode + retain it (R-ING-6),
    # exactly as extract_md does; otherwise src is resolved from disk later.
    data, ext = _decode_data_uri(src)
    doc.images.append(
        ImageRef(
            src=src,
            alt=node.attributes.get("alt") or "",
            caption=caption,
            heading_path=heading_path,
            data=data,
            ext=ext,
        )
    )


def _collect_links_images(node: Node, heading_path: str, doc: ExtractedDoc) -> None:
    for a in node.css("a[href]"):
        _add_link(a, doc)
    for img in node.css("img"):
        _add_image(img, heading_path, doc)


def _has_block_child(node: Node) -> bool:
    for child in node.iter(include_text=False):
        ct = child.tag
        if ct and ct.lower() in _BLOCK:
            return True
    return False


def _emit_body(doc: ExtractedDoc, text: str, heading_path: str) -> None:
    if text:
        doc.segments.append(Segment("body", text, heading_path))


def _walk(node: Node, stack: list[tuple[int, str]], doc: ExtractedDoc) -> None:
    tag = node.tag
    if not tag or tag.startswith("_") or tag == "-text":
        return
    tag = tag.lower()
    if tag in _SKIP:
        return

    if tag in _HEADINGS:
        text = _clean(node.text(separator=" "))
        _push_heading(stack, int(tag[1]), text)
        _emit_body(doc, text, _heading_path(stack))
        return

    if tag == "pre":  # code: preserve exact whitespace (no separator injection)
        code = node.text().rstrip("\n")
        if code.strip():
            doc.segments.append(Segment("code", code, _heading_path(stack)))
        _collect_links_images(node, _heading_path(stack), doc)
        return

    if tag == "table":
        text = _linearize_table(node)
        if text:
            doc.segments.append(Segment("table", text, _heading_path(stack)))
        _collect_links_images(node, _heading_path(stack), doc)
        return

    if tag == "figure":
        cap_node = node.css_first("figcaption")
        caption = _clean(cap_node.text(separator=" ")) if cap_node else ""
        for img in node.css("img"):
            _add_image(img, _heading_path(stack), doc, caption=caption)
        _emit_body(doc, caption, _heading_path(stack))
        return

    if tag == "img":
        _add_image(node, _heading_path(stack), doc)
        return

    # Leaf block: no block-level element child -> emit its whole text with a separator so
    # inline children (span/code/strong/a) never glue together (red-team Finding 1).
    if not _has_block_child(node):
        _emit_body(doc, _clean(node.text(separator=" ")), _heading_path(stack))
        _collect_links_images(node, _heading_path(stack), doc)
        return

    # Container with block children: emit its own inline/text runs, recurse block children.
    parts: list[str] = []

    def _flush() -> None:
        _emit_body(doc, _clean(" ".join(parts)), _heading_path(stack))
        parts.clear()

    for child in node.iter(include_text=True):
        ctag = child.tag
        if ctag == "-text":
            parts.append(child.text() or "")
        elif not ctag or ctag.startswith("_") or ctag.lower() in _SKIP:
            continue
        elif ctag.lower() in _BLOCK:
            _flush()
            _walk(child, stack, doc)
        else:  # inline element: keep its text with the current run, capture its links/images
            parts.append(child.text(separator=" "))
            if ctag.lower() == "a":
                _add_link(child, doc)
            _collect_links_images(child, _heading_path(stack), doc)
    _flush()


def extract_html(
    html: str,
    *,
    content_selector: str = "",
    strip_selectors: Sequence[str] = (),
) -> ExtractedDoc:
    """Extract structured content from one HTML document (§7.3).

    Applies ``strip_selectors`` then scopes to ``content_selector`` (falling back to the
    whole page, loudly, if the selector matches nothing). Preserves code blocks whole,
    linearizes tables with ``|``, and captures links + images with their heading path.
    """
    tree = HTMLParser(html)
    for selector in strip_selectors:
        for node in tree.css(selector):
            node.decompose()

    matched = True
    root: Node | None
    if content_selector:
        root = tree.css_first(content_selector)
        if root is None:
            matched = False
            root = tree.body
    else:
        root = tree.body
    if root is None:
        root = tree.root

    title = ""
    title_node = tree.css_first("title")
    if title_node is not None:
        title = _clean(title_node.text())
    if not title and root is not None:
        h1 = root.css_first("h1")
        if h1 is not None:
            title = _clean(h1.text())

    doc = ExtractedDoc(title=title, content_selector_matched=matched)
    if root is not None:
        _walk(root, [], doc)
    return doc


def _pdf_page_lines(page: object) -> list[tuple[float, bool, str]]:
    """One ``(font_size, bold, text)`` per visual line, in reading order — the font signal a PDF
    has instead of heading tags."""
    out: list[tuple[float, bool, str]] = []
    for block in page.get_text("dict").get("blocks", []):  # type: ignore[attr-defined]
        for line in block.get("lines", []):
            spans = line.get("spans", [])
            text = _clean("".join(str(s.get("text", "")) for s in spans))
            if not text:
                continue
            size = round(max((float(s.get("size", 0.0)) for s in spans), default=0.0), 1)
            bold = any(
                int(s.get("flags", 0)) & 16 or "bold" in str(s.get("font", "")).lower()
                for s in spans
            )
            out.append((size, bold, text))
    return out


def _pdf_heading_levels(pages: list[tuple[int, list[tuple[float, bool, str]]]]) -> dict[float, int]:
    """Map each heading-sized font to a level (1 = largest). The body size is the most common
    size (by character count); sizes clearly larger than it are headings, ranked descending.
    Returns {} when the document has no size variation (fall back to page locators)."""
    from collections import Counter

    weight: Counter[float] = Counter()
    for _, plines in pages:
        for size, _bold, text in plines:
            weight[size] += len(text)
    if not weight:
        return {}
    body_size = weight.most_common(1)[0][0]
    heading_sizes = sorted((s for s in weight if s > body_size * 1.12), reverse=True)
    return {size: min(idx + 1, 6) for idx, size in enumerate(heading_sizes)}


def extract_pdf(data: bytes) -> ExtractedDoc:
    """Extract text, **font-inferred heading structure**, links, and images from one PDF (§7.3).

    PyMuPDF is a ``[pdf]`` extra, imported lazily. A PDF has no heading tags, so headings are
    inferred from **font size/weight** (R-ING-4): the body size is the document's most common size,
    runs clearly larger than it become headings (ranked into levels), and body text is located by
    the resulting heading path (``Chapter > Section``) instead of a flat ``page N`` — so retrieval
    and citations point to the right section. A document with no size variation falls back to
    ``page N``. Link annotations → ``LinkRef``; embedded images are retained with inline bytes for
    the vision stage (R-ING-6).
    """
    import fitz  # lazy: the [pdf] extra is only loaded when a PDF is actually parsed

    doc_obj = fitz.open(stream=data, filetype="pdf")
    try:
        links: list[LinkRef] = []
        images: list[ImageRef] = []
        seen_xrefs: set[int] = set()  # a PDF image (xref) reused across pages is retained once
        pages: list[tuple[int, list[tuple[float, bool, str]]]] = []
        for i in range(doc_obj.page_count):
            page = doc_obj.load_page(i)
            pages.append((i, _pdf_page_lines(page)))
            locator = f"page {i + 1}"
            for link in page.get_links():
                uri = link.get("uri")
                if uri:
                    links.append(LinkRef(target=str(uri), anchor="", link_type="pdf_link"))
            for img_info in page.get_images(full=True):
                xref = int(img_info[0])
                if xref in seen_xrefs:
                    continue
                seen_xrefs.add(xref)
                try:
                    extracted = doc_obj.extract_image(xref)
                except Exception:  # noqa: BLE001 - a bad image must not abort the page
                    continue
                blob = extracted.get("image")
                if blob:
                    images.append(
                        ImageRef(src=f"pdf:xref{xref}", alt="", caption="", heading_path=locator,
                                 data=bytes(blob), ext=str(extracted.get("ext") or "png"))
                    )

        level_of = _pdf_heading_levels(pages)
        segments: list[Segment] = []
        stack: list[tuple[int, str]] = []
        buf: list[str] = []
        buf_loc: str | None = None

        def flush() -> None:
            nonlocal buf, buf_loc
            if buf:
                segments.append(Segment("body", " ".join(buf), buf_loc or ""))
                buf = []

        for i, plines in pages:
            page_loc = f"page {i + 1}"
            for size, _bold, text in plines:
                if size in level_of and len(text) <= 200:  # a heading (short, larger font)
                    if stack and stack[-1][1] == text:
                        continue  # same text at another size (e.g. title + H1) — not a new level
                    flush()
                    _push_heading(stack, level_of[size], text)
                    segments.append(Segment("body", text, _heading_path(stack)))
                else:  # body — located by the current heading path, else the page
                    loc = _heading_path(stack) or page_loc
                    if buf_loc is not None and buf_loc != loc:
                        flush()
                    buf_loc = loc
                    buf.append(text)
            flush()  # page boundary: keeps page-N-located body per page
            buf_loc = None

        title = _clean(str((doc_obj.metadata or {}).get("title") or ""))
        if not title and segments:
            title = _clean(segments[0].text.splitlines()[0])[:200]
    finally:
        doc_obj.close()
    return ExtractedDoc(title=title, segments=segments, links=links, images=images)


def _docx_heading_level(style_name: str) -> int:
    """`Heading 3` -> 3, `Title` -> 1, anything else -> 1 (a lone top-level heading)."""
    tail = style_name.rsplit(" ", 1)[-1]
    return int(tail) if tail.isdigit() else 1


def _docx_cell_text(cell: object, links: list[LinkRef]) -> str:
    """A cell's text = its own paragraphs **plus any nested tables** (recursed), while collecting
    hyperlinks found inside the cell (R-ING-5). ``_Cell.text`` alone drops nested tables, so we
    walk paragraphs and ``cell.tables`` explicitly — otherwise a table-in-a-cell is silent loss."""
    parts: list[str] = []
    for para in cell.paragraphs:  # type: ignore[attr-defined]
        t = _clean(para.text)
        if t:
            parts.append(t)
        for hl in para.hyperlinks:
            if hl.address:
                links.append(
                    LinkRef(target=hl.address, anchor=_clean(hl.text), link_type="docx_hyperlink")
                )
    for nested in cell.tables:  # type: ignore[attr-defined]
        nested_text = _linearize_docx_table(nested, links)
        if nested_text:
            parts.append(nested_text)
    return " ".join(parts)


def _linearize_docx_table(table: object, links: list[LinkRef]) -> str:
    """Flatten a python-docx table to the same ``|``-separated form HTML tables use (§7.3),
    recursing into nested tables and collecting in-cell hyperlinks (R-ING-5/6)."""
    rows: list[str] = []
    for row in table.rows:  # type: ignore[attr-defined]
        cells = [_docx_cell_text(cell, links) for cell in row.cells]
        if any(cells):
            rows.append(" | ".join(cells))
    return "\n".join(rows)


def _docx_para_images(para: object, heading_path: str, doc_obj: object, images: list[ImageRef]) -> None:
    """Retain inline images embedded in a paragraph's drawings (R-ING-6): pull each image part's
    bytes (carried inline for the vision stage) with its alt text (``wp:docPr`` descr/title)."""
    from docx.oxml.ns import qn

    for drawing in para._p.findall(".//" + qn("w:drawing")):  # type: ignore[attr-defined]
        docpr = drawing.find(".//" + qn("wp:docPr"))
        alt = _clean(docpr.get("descr") or docpr.get("title") or "") if docpr is not None else ""
        blip = drawing.find(".//" + qn("a:blip"))
        rid = blip.get(qn("r:embed")) if blip is not None else None
        if not rid:
            continue
        try:
            part = doc_obj.part.related_parts[rid]  # type: ignore[attr-defined]
        except (KeyError, AttributeError):
            continue
        blob = getattr(part, "blob", None)
        if not blob:
            continue
        ext = (Path(str(part.partname)).suffix or ".png").lstrip(".").lower()
        images.append(
            ImageRef(src=str(rid), alt=alt, caption="", heading_path=heading_path,
                     data=bytes(blob), ext=ext)
        )


def extract_docx(data: bytes) -> ExtractedDoc:
    """Extract headings, paragraphs, tables, hyperlinks, and inline images from one DOCX
    (§7.3, Phase 4b).

    python-docx is a ``[docx]`` extra, imported lazily so non-DOCX users never load it. Body
    elements are walked **in document order** (paragraphs and tables interleave) so heading paths
    and chunk order match the source: ``Heading N`` styles drive the heading path (R-ING-4),
    other paragraphs are ``body``, tables are ``table`` (linearized), w:hyperlink relationships
    become ``LinkRef``s (R-ING-5), and inline images are retained with their bytes for the vision
    stage (R-ING-6). Every extractor returns the common :class:`ExtractedDoc`.
    """
    from docx import Document  # lazy: the [docx] extra is only loaded when a DOCX is parsed
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc_obj = Document(io.BytesIO(data))
    segments: list[Segment] = []
    links: list[LinkRef] = []
    images: list[ImageRef] = []
    stack: list[tuple[int, str]] = []
    for child in doc_obj.element.body.iterchildren():
        if child.tag == qn("w:p"):
            para = Paragraph(child, doc_obj)
            text = _clean(para.text)
            style = (para.style.name if para.style else "") or ""
            if style.startswith(("Heading", "Title")):
                _push_heading(stack, _docx_heading_level(style), text)
                if text:
                    segments.append(Segment("body", text, _heading_path(stack)))
            elif text:
                segments.append(Segment("body", text, _heading_path(stack)))
            _docx_para_images(para, _heading_path(stack), doc_obj, images)
            for hl in para.hyperlinks:
                if hl.address:
                    links.append(
                        LinkRef(target=hl.address, anchor=_clean(hl.text), link_type="docx_hyperlink")
                    )
        elif child.tag == qn("w:tbl"):
            table_text = _linearize_docx_table(Table(child, doc_obj), links)
            if table_text:
                segments.append(Segment("table", table_text, _heading_path(stack)))
    title = _clean(str(doc_obj.core_properties.title or ""))
    if not title and segments:
        title = _clean(segments[0].text.splitlines()[0])[:200]
    return ExtractedDoc(title=title, segments=segments, links=links, images=images)


def _strip_front_matter(text: str) -> str:
    """Drop a leading YAML front-matter block (``---`` … ``---``, common in real Markdown like
    MDN) so its keys aren't indexed as body prose. An *unclosed* ``---`` block (no terminator in
    the first 64 lines) is treated as ordinary content, not silently swallowed."""
    if not (text.startswith("---\n") or text.startswith("---\r\n")):
        return text
    end = text.find("\n---", 3)
    # only strip a bounded, actually-closed block (guard against a stray leading '---')
    if end != -1 and text.count("\n", 0, end) <= 64:
        nl = text.find("\n", end + 1)
        return text[nl + 1 :] if nl != -1 else ""
    return text


def _md_inline_text(inline: object) -> str:
    """Rendered plain text of an ``inline`` token: text/code/image-alt children, breaks -> space
    (so ``[anchor](url)`` yields ``anchor``, not the raw markup)."""
    children = getattr(inline, "children", None)
    if not children:
        return _clean(getattr(inline, "content", "") or "")
    parts: list[str] = []
    for c in children:
        if c.type in ("text", "code_inline", "image"):  # image.content == alt text
            parts.append(c.content)
        elif c.type in ("softbreak", "hardbreak"):
            parts.append(" ")
    return _clean("".join(parts))


def _decode_data_uri(src: str) -> tuple[bytes | None, str]:
    """Decode a ``data:image/<ext>;base64,<data>`` URI to (bytes, ext), or (None, "") if it isn't
    one — so an image embedded inline in Markdown is retained (R-ING-6)."""
    import base64

    if not src.startswith("data:image/"):
        return None, ""
    header, _, payload = src.partition(",")
    if ";base64" not in header or not payload:
        return None, ""
    ext = header[len("data:image/") :].split(";", 1)[0].lower()
    try:
        # validate=True rejects non-base64 garbage (else b64decode silently keeps junk bytes)
        return base64.b64decode(payload, validate=True), ("jpg" if ext == "jpeg" else ext)
    except Exception:  # noqa: BLE001 - a malformed data URI just isn't retained
        return None, ""


def _collect_md_inline(
    inline: object, links: list[LinkRef], images: list[ImageRef], heading_path: str
) -> None:
    for c in getattr(inline, "children", None) or []:
        if c.type == "link_open":
            href = c.attrGet("href")
            if href:
                links.append(LinkRef(target=str(href), anchor="", link_type="md_link"))
        elif c.type == "image":
            src = str(c.attrGet("src") or "")
            if not src:
                continue
            data, ext = _decode_data_uri(src)
            images.append(
                ImageRef(
                    src="md-inline-image" if data is not None else src,
                    alt=_clean(c.content),
                    caption="",
                    heading_path=heading_path,
                    data=data,
                    ext=ext,
                )
            )


def _linearize_md_table(tokens: list, start: int, end: int) -> str:  # type: ignore[type-arg]
    """Flatten a markdown-it table token span to the ``a | b`` / newline-per-row form (§7.3)."""
    rows: list[str] = []
    cur: list[str] = []
    i = start
    while i <= end:
        t = tokens[i]
        if t.type == "tr_open":
            cur = []
        elif t.type in ("th_open", "td_open"):
            nxt = tokens[i + 1] if i + 1 <= end and tokens[i + 1].type == "inline" else None
            cur.append(_md_inline_text(nxt))
        elif t.type == "tr_close" and any(c.strip() for c in cur):
            rows.append(" | ".join(cur))
        i += 1
    return "\n".join(rows)


def extract_md(data: bytes) -> ExtractedDoc:
    """Extract headings, prose, fenced code, tables, links, and images from Markdown (§7.3,
    Phase 4c).

    markdown-it-py (the ``[md]`` extra) parses to a CommonMark token stream, imported lazily so
    non-MD users never load it. ``Heading N`` levels drive the heading path (R-ING-4), paragraphs
    are ``body``, fenced code is ``code`` (kept whole), GFM tables are ``table`` (linearized), and
    links/images become ``LinkRef``/``ImageRef`` (R-ING-5/6). Front matter is stripped. Every
    extractor returns the common :class:`ExtractedDoc`.
    """
    from markdown_it import MarkdownIt  # lazy: the [md] extra is only loaded to parse Markdown

    text = _strip_front_matter(data.decode("utf-8", errors="replace"))
    tokens = MarkdownIt("commonmark", {"html": False}).enable("table").parse(text)
    segments: list[Segment] = []
    links: list[LinkRef] = []
    images: list[ImageRef] = []
    stack: list[tuple[int, str]] = []
    n = len(tokens)
    i = 0
    while i < n:
        t = tokens[i]
        if t.type == "heading_open":
            inline = tokens[i + 1] if i + 1 < n else None
            htext = _md_inline_text(inline)
            _push_heading(stack, int(t.tag[1]), htext)
            if htext:
                segments.append(Segment("body", htext, _heading_path(stack)))
            _collect_md_inline(inline, links, images, _heading_path(stack))
            i += 3  # heading_open, inline, heading_close
            continue
        if t.type == "paragraph_open":
            inline = tokens[i + 1] if i + 1 < n else None
            ptext = _md_inline_text(inline)
            if ptext:
                segments.append(Segment("body", ptext, _heading_path(stack)))
            _collect_md_inline(inline, links, images, _heading_path(stack))
            i += 3
            continue
        if t.type == "fence":
            code = t.content.rstrip("\n")
            if code.strip():
                segments.append(Segment("code", code, _heading_path(stack)))
            i += 1
            continue
        if t.type == "table_open":
            j = i
            while j < n and tokens[j].type != "table_close":
                j += 1
            table_text = _linearize_md_table(tokens, i, j)
            if table_text:
                segments.append(Segment("table", table_text, _heading_path(stack)))
            i = j + 1
            continue
        i += 1
    title = _clean(segments[0].text.splitlines()[0])[:200] if segments else ""
    return ExtractedDoc(title=title, segments=segments, links=links, images=images)


def _linearize_pptx_table(table: object) -> str:
    """Linearize a PPTX table to pipe-delimited rows (matching the HTML/DOCX table shape)."""
    rows: list[str] = []
    for row in table.rows:  # type: ignore[attr-defined]
        cells = [_clean(cell.text) for cell in row.cells]
        if any(cells):
            rows.append(" | ".join(cells))
    return "\n".join(rows)


def _pptx_add_image(shape: object, heading_path: str, images: list[ImageRef]) -> None:
    """Retain a picture shape's bytes (R-ING-6). Alt text comes from the shape's descr / name."""
    try:
        img = shape.image  # type: ignore[attr-defined]
    except Exception:  # noqa: BLE001 - non-picture or unreadable blob: skip, don't abort the slide
        return
    alt = ""
    with contextlib.suppress(Exception):
        descr = shape._element.xpath(".//p:cNvPr/@descr")  # type: ignore[attr-defined] # noqa: SLF001
        alt = _clean(descr[0]) if descr else _clean(getattr(shape, "name", "") or "")
    images.append(
        ImageRef(src="", alt=alt, caption="", heading_path=heading_path,
                 data=img.blob, ext=(img.ext or "png")),
    )


def extract_pptx(data: bytes) -> ExtractedDoc:
    """Extract slides from a PPTX (§7.3, Phase 6 / GATE 6): each slide's title → a level-1 heading,
    body text frames → ``body``, **speaker notes** → ``body``, tables → ``table`` (linearized), and
    pictures retained with their bytes (R-ING-6). The slide number anchors the heading path so a
    citation points at the right slide. python-pptx is the ``[pptx]`` extra, imported lazily."""
    from pptx import Presentation  # lazy: only when a PPTX is parsed
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(data))
    segments: list[Segment] = []
    images: list[ImageRef] = []
    for idx, slide in enumerate(prs.slides, start=1):
        title_shape = slide.shapes.title
        title = _clean(title_shape.text) if title_shape is not None else ""
        hp = _heading_path([(1, f"Slide {idx}: {title}" if title else f"Slide {idx}")])
        if title:
            segments.append(Segment("body", title, hp))
        for shape in slide.shapes:
            if title_shape is not None and shape == title_shape:
                continue
            if getattr(shape, "has_table", False):
                table_text = _linearize_pptx_table(shape.table)
                if table_text:
                    segments.append(Segment("table", table_text, hp))
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                _pptx_add_image(shape, hp, images)
            elif getattr(shape, "has_text_frame", False):
                text = _clean(shape.text_frame.text)
                if text:
                    segments.append(Segment("body", text, hp))
        if slide.has_notes_slide:
            notes = _clean(slide.notes_slide.notes_text_frame.text)
            if notes:
                nhp = _heading_path([(1, f"Slide {idx}: {title}" if title else f"Slide {idx}"),
                                     (2, "Speaker notes")])
                segments.append(Segment("body", notes, nhp))
    title = _clean(str(prs.core_properties.title or ""))
    if not title and segments:
        title = _clean(segments[0].text.splitlines()[0])[:200]
    return ExtractedDoc(title=title, segments=segments, images=images)


def extract_document(
    path: Path,
    ext: str,
    *,
    content_selector: str = "",
    strip_selectors: Sequence[str] = (),
) -> ExtractedDoc:
    """Dispatch to the right extractor by file extension — the pluggable-parser seam (R-PROC-6).

    HTML via selectolax (content_selector/strip_selectors apply), PDF via PyMuPDF, DOCX via
    python-docx, Markdown via markdown-it-py. A new format adds one branch here plus its
    extractor; the rest of the pipeline (chunker, links, store) is format-agnostic because every
    extractor returns an :class:`ExtractedDoc`.
    """
    fmt = ext.lower().lstrip(".")
    if fmt == "pdf":
        return extract_pdf(path.read_bytes())
    if fmt == "docx":
        return extract_docx(path.read_bytes())
    if fmt in ("md", "markdown"):
        return extract_md(path.read_bytes())
    if fmt == "pptx":
        return extract_pptx(path.read_bytes())
    html = path.read_bytes().decode("utf-8", errors="replace")
    return extract_html(
        html, content_selector=content_selector, strip_selectors=list(strip_selectors)
    )


# ----------------------------------------------------------------- chunking


@dataclass
class Chunk:
    """A ready-to-index unit of a document, with its heading-path locator (R-ING-4)."""

    ord: int
    kind: str  # body | code
    locator: str
    text: str


def chunk_document(doc: ExtractedDoc, *, chunk_tokens: int, overlap: int) -> list[Chunk]:
    """Split a document into chunks (§7.6).

    Body/table text is grouped by heading path and packed toward ``chunk_tokens`` words
    with a sliding ``overlap``; a **code** segment always becomes one whole chunk, never
    split. Each chunk records its heading-path locator (R-ING-4).
    """
    ov = min(max(overlap, 0), max(chunk_tokens - 1, 0))
    chunks: list[Chunk] = []
    buf: list[str] = []
    hp: str | None = None
    emitted = False  # has this heading group produced a chunk yet?
    fresh = False  # has new (non-overlap) content been added since the last emit?

    def emit() -> None:
        nonlocal emitted
        if buf:
            chunks.append(Chunk(len(chunks), "body", hp or "", " ".join(buf)))
            emitted = True

    def close_group() -> None:
        nonlocal buf, fresh
        if buf and (fresh or not emitted):
            emit()
        buf = []
        fresh = False

    for seg in doc.segments:
        if seg.kind == "code":
            close_group()
            emitted = False
            chunks.append(Chunk(len(chunks), "code", seg.heading_path, seg.text))
            hp = None
            continue

        if hp is None:
            hp, emitted, fresh = seg.heading_path, False, False
        elif seg.heading_path != hp:
            close_group()
            hp, emitted = seg.heading_path, False

        for word in seg.text.split():
            buf.append(word)
            fresh = True
            if len(buf) >= chunk_tokens:
                emit()
                buf = buf[-ov:] if ov > 0 else []
                fresh = False

    close_group()
    return chunks


# ----------------------------------------------------------------- orchestration

_EXTERNAL_PREFIXES = ("http://", "https://", "ftp://", "mailto:", "javascript:", "tel:", "data:")


@dataclass
class IngestResult:
    """Counts and audit data for one ingest run (§7.8) — the input to Gate 1."""

    files_found: int = 0
    included: int = 0
    excluded_glob: int = 0
    other_files: int = 0
    skipped_unchanged: int = 0
    stripped_empty: int = 0
    parse_errors: int = 0
    documents: int = 0
    chunks: int = 0
    embedded: int = 0
    images: int = 0
    images_unresolved: int = 0  # referenced but not retained (missing on disk / remote)
    gotchas: int = 0  # chunks dual-tagged [GOTCHA] + flags row (R-ING-8)
    relations_total: int = 0
    relations_resolved: int = 0
    relations_unresolved: int = 0
    content_selector_misses: int = 0
    zero_chunk_docs: int = 0
    untagged_audience_docs: int = 0
    per_extension: dict[str, int] = field(default_factory=dict)
    errors: list[tuple[str, str]] = field(default_factory=list)
    timings_ms: dict[str, float] = field(default_factory=dict)


def _resolve_local(raw: str, src_path: Path) -> Path | None:
    """Resolve a raw href/src to a local path, or None if it is an external/data URL."""
    target = raw.split("#", 1)[0].split("?", 1)[0].strip()
    if not target or target.startswith("//") or target.lower().startswith(_EXTERNAL_PREFIXES):
        return None
    return (src_path.parent / target).resolve()


def _stage_images(
    doc: ExtractedDoc,
    src_path: Path,
    staging_dir: Path,
    store: Store,
    doc_id: int,
    base_ord: int,
    result: IngestResult,
) -> int:
    """Retain image originals (keyed by sha256) and add searchable image_ref chunks."""
    images_dir = staging_dir / "images"
    added = 0
    for img in doc.images:
        if img.data is not None:  # inline image bytes (PDF/DOCX embedded image)
            data: bytes | None = img.data
            ext = (img.ext or "bin").lstrip(".").lower()
        else:  # HTML: resolve the referenced file on disk
            local = _resolve_local(img.src, src_path)
            if local is not None and local.is_file():
                data = local.read_bytes()
                ext = local.suffix.lstrip(".").lower() or "bin"
            else:
                data = None
                if local is not None:  # a LOCAL ref that didn't resolve to a file — actionable,
                    result.images_unresolved += 1  # surfaced in the audit (not silently dropped)
        if data is not None:
            sha = hashlib.sha256(data).hexdigest()
            images_dir.mkdir(parents=True, exist_ok=True)
            dest = images_dir / f"{sha}.{ext}"
            if not dest.exists():
                dest.write_bytes(data)
            store.add_image(
                sha256=sha,
                ext=ext,
                doc_id=doc_id,
                locator=img.heading_path,
                alt=img.alt,
                caption=img.caption,
                num_bytes=len(data),
            )
            result.images += 1
        text = " ".join(t for t in (img.alt, img.caption, img.heading_path) if t).strip()
        if text:
            store.add_chunk(
                document_id=doc_id,
                ord=base_ord + added,
                text=text,
                kind="image_ref",
                locator=img.heading_path,
            )
            added += 1
    return added


# --- parallel parse (§7) -----------------------------------------------------------
# The CPU-heavy part of ingest (read + extract + chunk) runs in worker processes; the DB
# writes stay in the single main-process writer and are applied in file order, so the
# result is byte-identical to a serial run (R-SRCH-5) — parallelism only speeds it up.


@dataclass
class _ParseTask:
    """One file's parse job — picklable, so it crosses the spawn boundary to a worker."""

    path: str  # resolved posix path (also the document's DB key)
    ext: str
    content_selector: str
    strip_selectors: tuple[str, ...]
    chunk_tokens: int
    chunk_overlap: int
    stored_hash: str | None  # the DB's current hash for this path (for the skip decision)
    force: bool


@dataclass
class _ParseResult:
    """A worker's output for one file — picklable, consumed serially in the main process."""

    path: str
    ext: str
    status: str  # "skip" | "ok" | "error"
    file_hash: str
    mtime: float
    error: str | None
    doc: ExtractedDoc | None
    chunks: list[Chunk]


def _parse_file(task: _ParseTask) -> _ParseResult:
    """Worker: hash, extract, and chunk one file. Pure CPU/IO, touches no DB — this is the
    part that runs in a process pool. Unchanged files short-circuit before the parse."""
    path = Path(task.path)
    try:
        file_hash = content_hash(path)
    except OSError as err:
        return _ParseResult(
            task.path, task.ext, "error", "", 0.0, f"{type(err).__name__}: {err}", None, []
        )
    if task.stored_hash is not None and task.stored_hash == file_hash and not task.force:
        return _ParseResult(task.path, task.ext, "skip", file_hash, 0.0, None, None, [])
    try:
        doc = extract_document(
            path,
            task.ext,
            content_selector=task.content_selector,
            strip_selectors=task.strip_selectors,
        )
        chunks = chunk_document(doc, chunk_tokens=task.chunk_tokens, overlap=task.chunk_overlap)
    except Exception as err:  # noqa: BLE001 - a bad file is reported, not fatal
        return _ParseResult(
            task.path, task.ext, "error", file_hash, 0.0, f"{type(err).__name__}: {err}", None, []
        )
    return _ParseResult(
        task.path, task.ext, "ok", file_hash, path.stat().st_mtime, None, doc, chunks
    )


def _tag_gotchas(
    chunks: Sequence[Chunk], patterns: Sequence[GotchaPattern]
) -> tuple[list[tuple[int, str, str, str]], list[tuple[int, str]]]:
    """Apply approved gotcha rules to a document's chunks (R-ING-8). Returns the ``add_chunks``
    tuples (text prefixed with ``[GOTCHA]`` where a rule matched) plus ``(ord, label)`` for each
    match so the caller can write the paired ``flags`` row. No patterns ⇒ chunks unchanged."""
    tuples: list[tuple[int, str, str, str]] = []
    matched: list[tuple[int, str]] = []
    for c in chunks:
        text = c.text
        if patterns:
            label = match_gotcha(text, list(patterns))
            if label is not None:
                text = gotcha_tag_text(text)
                matched.append((c.ord, label))
        tuples.append((c.ord, text, c.kind, c.locator))
    return tuples, matched


def _write_parsed(
    res: _ParseResult,
    source: SourceConfig,
    config: Config,
    store: Store,
    result: IngestResult,
    gotcha_patterns: Sequence[GotchaPattern] = (),
) -> None:
    """Main process: turn one parse result into DB rows (single writer, called in file order
    so doc/chunk ids are assigned deterministically)."""
    result.per_extension[res.ext] = result.per_extension.get(res.ext, 0) + 1
    if res.status == "skip":
        result.skipped_unchanged += 1
        return
    if res.status == "error" or res.doc is None:
        result.parse_errors += 1
        result.errors.append((res.path, res.error or "parse failed"))
        return
    doc = res.doc
    old = store.document_id_for_path(res.path)  # re-ingest of a changed/forced file
    if old is not None:
        store.delete_document(old)
    if not doc.content_selector_matched:
        result.content_selector_misses += 1
    if doc.text_length < source.min_content_chars:
        result.stripped_empty += 1
        return
    if not source.audience:
        result.untagged_audience_docs += 1
    doc_id = store.add_document(
        path=res.path,
        source=source.name,
        source_version=source.version,
        title=doc.title,
        content_hash=res.file_hash,
        content_type="documentation",
        fmt=res.ext,
        audience=source.audience,
        mtime=res.mtime,
        status="active",
    )
    result.documents += 1
    tagged, matched = _tag_gotchas(res.chunks, gotcha_patterns)
    store.add_chunks(doc_id, tagged)
    if matched:  # dual-tag: text prefix (done above) + a filterable flags row (R-ING-8)
        by_ord = {int(r["ord"]): int(r["id"]) for r in store.chunks_for_document(doc_id)}
        for ordv, label in matched:
            store.add_flag(
                doc_id=doc_id, chunk_id=by_ord.get(ordv), kind="gotcha",
                source="preflight", rule_id=label,
            )
        result.gotchas += len(matched)
    image_chunks = _stage_images(
        doc, Path(res.path), Path(config.paths.staging_dir), store, doc_id, len(res.chunks), result
    )
    total_chunks = len(res.chunks) + image_chunks
    result.chunks += total_chunks
    if total_chunks == 0:
        result.zero_chunk_docs += 1
    for link in doc.links:
        store.add_relation(src_doc=doc_id, dst_raw=link.target, link_type=link.link_type)


def _resolve_workers(requested: int | None, n_tasks: int) -> int:
    """How many parse workers to use. Auto: parallelize only when there's enough work to
    amortize process-spawn cost. ``DOCUSEARCH_INGEST_WORKERS`` overrides; 1 = serial."""
    if requested is None:
        env = os.environ.get("DOCUSEARCH_INGEST_WORKERS")
        if env:
            try:
                requested = int(env)
            except ValueError:
                requested = None
    if requested is not None:
        return max(1, requested)
    if n_tasks < 250:  # small ingests: spawn overhead isn't worth it
        return 1
    return min(os.cpu_count() or 1, 8)


def _resolve_links(store: Store) -> None:
    """Post-pass: resolve raw link targets to document ids now that all docs exist (R-ING-5)."""
    path_to_id = store.document_path_to_id()
    for rel_id, src_path, dst_raw in store.unresolved_relations():
        local = _resolve_local(dst_raw, Path(src_path))
        if local is not None:
            dst = path_to_id.get(local.as_posix())
            if dst is not None:
                store.set_relation_dst(rel_id, dst)


def _embed_chunks(
    store: Store,
    provider: EmbedProvider,
    batch_size: int,
    *,
    progress: ProgressFn | None = None,
) -> int:
    """Batch-embed every chunk that lacks a vector and store it, tagged by model (§7.7).

    Only un-embedded chunks are processed, so incremental re-ingest embeds just the new
    chunks. Refuses to mix models in one index (R-EMB-2/3): a store already embedded by a
    different model must be re-indexed fresh.
    """
    # Provenance guard runs FIRST — before the no-pending early-out — so a swap to a
    # different model is refused even when no new chunks need embedding (R-EMB-3: never
    # silently mix models in one index). We read the *embeddings rows*, not the
    # ``embed_model`` meta flag: the flag is written only after a full pass, so an
    # interrupted run leaves committed vectors with no flag — exactly the gap that used
    # to let a later, different-dimension model mix in and crash the ANN build.
    existing = store.existing_embedding_model()
    if existing is not None and existing[0] != provider.model_id:
        old_model, old_dim = existing
        raise embed.EmbedError(
            f"index already holds vectors from {old_model!r} (dim {old_dim}), but the "
            f"configured model is {provider.model_id!r}. Vectors from different models "
            f"can't share one index. Recover with either:\n"
            f"  - docusearch ingest --reembed   (drop old vectors, re-embed with the new model)\n"
            f"  - point paths.db_path at a fresh database\n"
            f"(this can happen if a previous embedding run was interrupted mid-way)."
        )
    pending = store.chunks_without_embeddings()
    if not pending:
        return 0
    dim = provider.dim  # loads the model (may download) — only reached when there is work
    if existing is not None and existing[1] != dim:  # same name, different dim: still refuse
        raise embed.EmbedError(
            f"index holds {existing[1]}-dim vectors but {provider.model_id!r} produces "
            f"{dim}-dim; run `docusearch ingest --reembed` or use a fresh db_path."
        )
    # Record provenance up front so an interrupted run stays self-describing (the guard
    # above then catches any later model swap instead of silently mixing).
    store.set_meta("embed_model", provider.model_id)
    store.set_meta("embed_dim", str(dim))
    total = len(pending)
    embedded = 0
    for start in range(0, total, batch_size):
        batch = pending[start : start + batch_size]
        vectors = provider.embed([text for _, text in batch])
        rows = [
            (chunk_id, provider.model_id, dim, embed.to_blob(vectors[i]))
            for i, (chunk_id, _) in enumerate(batch)
        ]
        store.add_embeddings(rows)
        embedded += len(rows)
        if progress is not None:
            progress("embed", embedded, total)
    return embedded


_FROM_CONFIG: EmbedProvider | None = object()  # type: ignore[assignment]


def run_ingest(
    config: Config,
    store: Store,
    *,
    force: bool = False,
    reembed: bool = False,
    provider: EmbedProvider | None = _FROM_CONFIG,
    progress: ProgressFn | None = None,
    workers: int | None = None,
) -> IngestResult:
    """Ingest every configured source into the store and return an audit result (§7).

    ``provider`` defaults to whatever ``embed.model`` selects (None => BM25-only,
    R-CFG-4); pass an explicit provider to override (used by tests/harness).
    ``force`` re-parses every file AND rebuilds all vectors (a full rebuild); ``reembed``
    rebuilds only the vectors (keeps parsed docs — the cheap model switch). Both drop
    existing vectors up front, which also heals an orphaned/mixed embeddings table.
    ``progress`` receives (phase, done, total) callbacks for a live bar (CLI only).
    ``workers`` sets the parse-pool size (None => auto by corpus size; 1 => serial;
    ``DOCUSEARCH_INGEST_WORKERS`` overrides). Results are always applied in file order,
    so the worker count never changes the index (R-SRCH-5).
    """
    prov = embed.make_provider(config.embed) if provider is _FROM_CONFIG else provider
    result = IngestResult()
    start = time.perf_counter()
    # Approved gotcha rules (R-ING-8) — empty unless preflight_rules.yaml exists AND approved.
    gotcha_patterns = active_gotcha_patterns(config.enrich.preflight_rules)
    if reembed or force:
        # A full rebuild (--force) and an explicit --reembed both start vectors from a
        # clean slate. This is authoritative: it clears vectors that the per-document
        # re-ingest cascade would miss — orphans, or docs whose path changed — instead of
        # relying on the cascade and then tripping the model-mismatch guard on stragglers.
        store.clear_embeddings()
    # Classify every source up front so the file total (for progress) is known before we
    # start the slow parse loop.
    worklist: list[tuple[Path, SourceConfig]] = []
    for source in config.sources:
        root = Path(source.location)
        if not root.is_dir():
            result.errors.append((source.location, "source location not found"))
            continue
        included, excluded_glob, other = _classify_files(root, source.include, source.exclude)
        result.files_found += len(included) + excluded_glob + other
        result.included += len(included)
        result.excluded_glob += excluded_glob
        result.other_files += other
        worklist += [(path, source) for path in included]

    total_files = len(worklist)
    # Build one picklable parse task per file (the DB's current hashes let workers decide
    # skip-vs-reparse without touching the DB).
    stored = store.document_hashes()
    tasks = [
        _ParseTask(
            path=path.resolve().as_posix(),
            ext=(path.suffix.lstrip(".").lower() or "html"),
            content_selector=source.content_selector,
            strip_selectors=tuple(source.strip_selectors),
            chunk_tokens=config.index.chunk_tokens,
            chunk_overlap=config.index.chunk_overlap,
            stored_hash=stored.get(path.resolve().as_posix()),
            force=force,
        )
        for path, source in worklist
    ]
    sources = [source for _, source in worklist]
    n_workers = _resolve_workers(workers, total_files)

    # Parse in parallel (CPU-bound extract+chunk across cores — the real fix for the slow,
    # ~1-core-bound ingest, R-PERF-1), but apply DB writes serially IN FILE ORDER so ids are
    # assigned deterministically and the result is byte-identical to a serial run (R-SRCH-5).
    # Commits are batched per file (deferred_commits) to avoid per-insert fsync churn.
    with store.deferred_commits():
        if n_workers <= 1:
            parsed: Iterator[_ParseResult] = (_parse_file(t) for t in tasks)
            for done, (res, source) in enumerate(zip(parsed, sources, strict=True), 1):
                _write_parsed(res, source, config, store, result, gotcha_patterns)
                store.commit()
                if progress is not None:
                    progress("ingest", done, total_files)
        else:
            import concurrent.futures as _cf
            import multiprocessing as _mp

            ctx = _mp.get_context("spawn")  # spawn-safe: workers re-import the module
            with _cf.ProcessPoolExecutor(max_workers=n_workers, mp_context=ctx) as pool:
                ordered = pool.map(_parse_file, tasks, chunksize=8)  # results stay in order
                for done, (res, source) in enumerate(zip(ordered, sources, strict=True), 1):
                    _write_parsed(res, source, config, store, result, gotcha_patterns)
                    store.commit()
                    if progress is not None:
                        progress("ingest", done, total_files)

    with store.deferred_commits():  # link resolution is another burst of small updates
        _resolve_links(store)
    result.relations_total = store.count_relations()
    result.relations_resolved = store.count_resolved_relations()
    result.relations_unresolved = result.relations_total - result.relations_resolved

    if prov is not None:  # embed at index time unless embed.model: none (R-CFG-4)
        embed_start = time.perf_counter()
        result.embedded = _embed_chunks(store, prov, config.embed.batch_size, progress=progress)
        result.timings_ms["embed"] = (time.perf_counter() - embed_start) * 1000.0
        # Build/refresh the ANN sidecar for file-backed indexes (§7.7). In-memory stores
        # fall back to numpy brute-force at query time, so no sidecar is needed there.
        db_path = config.paths.db_path
        if config.index.ann and db_path != ":memory:" and store.count_embeddings() > 0:
            from . import search

            search.VectorIndex.build(
                store,
                prov.dim,
                Path(db_path).with_suffix(".hnsw"),
                m=config.index.ann_m,
                ef_construction=config.index.ann_ef_construction,
            )

    result.timings_ms["total"] = (time.perf_counter() - start) * 1000.0
    runlog.log(
        "ingest.done",
        documents=result.documents,
        chunks=result.chunks,
        embedded=result.embedded,
        images=result.images,
        skipped_unchanged=result.skipped_unchanged,
        stripped_empty=result.stripped_empty,
        parse_errors=result.parse_errors,
        relations_unresolved=result.relations_unresolved,
        took_ms=round(result.timings_ms["total"], 1),
    )
    return result


# ----------------------------------------------------------------- audit rendering


def render_ingest_audit(result: IngestResult, *, run_id: str = "") -> str:
    """Render the ingest audit report (§7.8) — the input to Gate 1.

    Loudly surfaces everything that was skipped or is suspect (nothing dropped
    silently): glob exclusions, too-short strips, parse errors, content-selector
    misses, zero-chunk docs, untagged audiences, and unresolved links.
    """
    took = result.timings_ms.get("total", 0.0)
    lines = [
        "# Ingest audit",
        "",
        f"run_id: `{run_id}`  ·  elapsed: {took / 1000:.2f}s",
        "",
        "## Discovery",
        f"- files found: **{result.files_found}**",
        f"- included (passed globs): **{result.included}**",
        f"- excluded by glob: **{result.excluded_glob}**",
        f"- other (present, not selected): **{result.other_files}**",
        "",
        "## Documents & chunks",
        f"- documents ingested this run: **{result.documents}**",
        f"- skipped (unchanged hash): **{result.skipped_unchanged}**",
        f"- chunks written: **{result.chunks}**",
        f"- chunks embedded: **{result.embedded}**",
        f"- images retained: **{result.images}**",
        f"- gotchas dual-tagged (approved rules): **{result.gotchas}**",
        "",
        "## Relations (cross-references)",
        f"- total: **{result.relations_total}**",
        f"- resolved: **{result.relations_resolved}**",
        f"- unresolved (external / broken — kept for audit): **{result.relations_unresolved}**",
        "",
        "## ⚠️ Skips & anomalies (review every non-zero line)",
        f"- stripped as too short (< min_content_chars): **{result.stripped_empty}**",
        f"- image refs unresolved (local file missing — not retained/enrichable): "
        f"**{result.images_unresolved}**",
        f"- content_selector matched nothing (fell back to body): **{result.content_selector_misses}**",
        f"- zero-chunk documents: **{result.zero_chunk_docs}**",
        f"- documents with no audience tag: **{result.untagged_audience_docs}**",
        f"- parse errors: **{result.parse_errors}**",
        "",
        "## Per-extension",
    ]
    if result.per_extension:
        lines += [f"- `{ext}`: {n}" for ext, n in sorted(result.per_extension.items())]
    else:
        lines.append("- (none)")
    if result.errors:
        lines += ["", "## Errors", *[f"- `{path}` — {msg}" for path, msg in result.errors[:200]]]
        if len(result.errors) > 200:
            lines.append(f"- … and {len(result.errors) - 200} more")
    lines.append("")
    return "\n".join(lines)


def render_store_audit(store: Store) -> str:
    """Render the current index state for ``docusearch audit`` (spot-check counts)."""
    resolved = store.count_resolved_relations()
    total_rel = store.count_relations()
    histogram = store.fmt_histogram()
    fmt_lines = [f"- `{fmt}`: {n}" for fmt, n in histogram.items()] or ["- (none)"]
    lines = [
        "# Index audit",
        "",
        f"- documents: **{store.count_documents()}**",
        f"- chunks: **{store.count_chunks()}**",
        f"- images: **{store.count_images()}**",
        f"- gotcha flags: **{store.count_flags('gotcha')}**",
        f"- relations: **{total_rel}** ({resolved} resolved, {total_rel - resolved} unresolved)",
        "",
        "## By format",
        *fmt_lines,
        "",
        "## ⚠️ Anomalies",
        f"- documents with zero chunks: **{store.documents_without_chunks()}**",
        f"- documents with empty audience: **{store.documents_with_empty_audience()}**",
        "",
    ]
    return "\n".join(lines)
